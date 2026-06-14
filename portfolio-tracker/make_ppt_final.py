from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, io
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── 색상 ──────────────────────────────────────────────────
INDIGO   = RGBColor(0x3F, 0x51, 0xB5)
INDIGO_D = RGBColor(0x1A, 0x23, 0x7E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x21, 0x21, 0x21)
GRAY     = RGBColor(0x9E, 0x9E, 0x9E)
LGRAY    = RGBColor(0xF5, 0xF5, 0xF5)
BORDER   = RGBColor(0xE0, 0xE0, 0xE0)
GREEN    = RGBColor(0x2E, 0x7D, 0x32)
AMBER    = RGBColor(0xE6, 0x5C, 0x00)

BLANK = prs.slide_layouts[6]
TOTAL_SLIDES = 19

# ── 헬퍼 ──────────────────────────────────────────────────
def rect(s, l, t, w, h, fill=None, line_color=None, lw=0.8):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line_color:
        sh.line.color.rgb = line_color; sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh

def txt(s, text, l, t, w, h, size=14, bold=False,
        color=DARK, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    return tb

def header(s, title, sub=None):
    rect(s, 0, 0, 13.33, 1.0)
    txt(s, title, 0.5, 0.08, 12, 0.62, size=24, bold=True, color=INDIGO)
    if sub:
        txt(s, sub, 0.5, 0.68, 12, 0.3, size=11, color=GRAY)
    rect(s, 0, 0.95, 13.33, 0.05, fill=INDIGO)

def pagenum(s, n):
    txt(s, f"{n} / {TOTAL_SLIDES}", 12.3, 7.1, 0.9, 0.3,
        size=10, color=GRAY, align=PP_ALIGN.CENTER)

def chip(s, label, l, t, w, h, selected=False):
    if selected:
        rect(s, l, t, w, h, fill=INDIGO)
        txt(s, "✓  " + label, l+0.15, t+0.12, w-0.2, h-0.15,
            size=16, bold=True, color=WHITE)
    else:
        rect(s, l, t, w, h, line_color=BORDER)
        txt(s, label, l+0.15, t+0.12, w-0.2, h-0.15,
            size=16, color=GRAY)


# ════════════════════════════════════════════════════════════
# 1. 표지
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=WHITE)
rect(s, 0, 5.1, 13.33, 0.05, fill=INDIGO)
rect(s, 0, 5.15, 13.33, 2.35, fill=LGRAY)

txt(s, "Portfolio Tracker", 1, 1.0, 11, 1.4,
    size=52, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
txt(s, "되돌아보는 투자자가 앞으로 나아간다",
    1, 2.6, 11, 0.6, size=20, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, "최종 발표",
    1, 3.4, 11, 0.55, size=18, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)

txt(s, "Flutter + Firebase  ·  개인 자산 포트폴리오 관리 앱",
    1, 5.5, 11, 0.5, size=14, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, "2026.06",
    1, 6.0, 11, 0.45, size=12, color=GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# 2. 문제 정의
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "1.  문제 정의", "누구의 어떤 문제를 푸는가")
pagenum(s, 2)

problems = [
    ("데이터 분산",  "여러 계좌·앱에 흩어진 투자 기록 — 전체 현황 파악 불가능"),
    ("복기 불가",    "체계적인 기록 수단 없음 — 매수 근거와 실제 결과를 나란히 볼 수 없음"),
    ("실수 반복",    "복기 부재 → 같은 판단 오류 반복 → 수익 기회 상실"),
]
for i, (label, desc) in enumerate(problems):
    ty = 1.15 + i * 1.5
    rect(s, 0.5, ty, 12.3, 1.25, line_color=BORDER)
    rect(s, 0.5, ty, 0.06, 1.25, fill=INDIGO)
    txt(s, label, 0.72, ty + 0.15, 2.8, 0.42, size=13, bold=True, color=INDIGO)
    txt(s, desc, 0.72, ty + 0.6, 11.3, 0.55, size=15, color=DARK)

rect(s, 0.5, 5.7, 12.3, 0.95, fill=LGRAY, line_color=BORDER)
rect(s, 0.5, 5.7, 0.06, 0.95, fill=INDIGO)
txt(s, '한 줄 가치 제안  ·  "경험을 자산으로 바꾸는 개인 투자 포트폴리오 관리 앱"',
    0.72, 5.9, 11.8, 0.55, size=14, bold=True, color=DARK)


# ════════════════════════════════════════════════════════════
# 3. 활용 방안 & 기대 효과
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "2.  활용 방안 & 기대 효과", "이 앱으로 무엇을 할 수 있는가")
pagenum(s, 3)

use_items = [
    ("투자 기록 통합 관리",
     "주식·코인 계좌를 한 앱에서 관리  ·  입출금·잔액·수익률 자동 계산\n"
     "여러 계좌에 흩어진 투자 현황을 대시보드 한 화면에서 즉시 파악"),
    ("복기를 통한 실력 향상",
     "매수 근거 · 목표가 · 손절가를 기록하고, 정리 후 실제 손익과 나란히 비교\n"
     "패턴 파악 → 반복 실수 방지 → 장기적으로 판단력 향상"),
    ("지속적인 성과 추적",
     "수익률 막대 차트 · 누적 손익 라인 차트 · 기간 필터로 성과 시각화\n"
     "승률·누적 손익 등 핵심 지표를 대시보드에서 상시 모니터링"),
]
for i, (title, body) in enumerate(use_items):
    ty = 1.15 + i * 1.82
    rect(s, 0.5, ty, 12.3, 1.55, line_color=BORDER)
    rect(s, 0.5, ty, 0.06, 1.55, fill=GREEN)
    rect(s, 0.5, ty, 12.3, 0.48, fill=LGRAY)
    txt(s, title, 0.72, ty + 0.1, 12.0, 0.3, size=13, bold=True, color=INDIGO)
    txt(s, body,  0.72, ty + 0.56, 11.8, 0.88, size=12, color=DARK)


# ════════════════════════════════════════════════════════════
# 4. 사용자 시나리오
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "3.  사용자 시나리오", "포지션 기록과 복기 — 핵심 유즈케이스")
pagenum(s, 4)

# 시나리오 텍스트
rect(s, 0.5, 1.15, 12.3, 1.6, fill=LGRAY, line_color=BORDER)
rect(s, 0.5, 1.15, 0.06, 1.6, fill=INDIGO)
txt(s, "Min은 삼성전자를 68,000원에 100주 매수했다.",
    0.72, 1.25, 12.0, 0.45, size=15, bold=True, color=DARK)
txt(s, "앱에 진입가 · 수량 · 매수 근거 · 목표가 · 손절가를 기록하고,",
    0.72, 1.7, 12.0, 0.4, size=14, color=DARK)
txt(s, "72,000원에 매도하면 수익률(+5.9%) · 손익(+400,000원)이 자동 계산된다.",
    0.72, 2.1, 12.0, 0.4, size=14, color=DARK)

# 흐름 다이어그램
steps = [
    ("1", "로그인"),
    ("2", "계좌\n등록"),
    ("3", "포지션\n진입 기록"),
    ("4", "포지션\n정리 기록"),
    ("5", "복기\n메모"),
]
sw, sh_box, gap = 2.2, 1.6, 0.08
start_x = (13.33 - (len(steps) * sw + (len(steps)-1) * gap)) / 2

for i, (n, label) in enumerate(steps):
    sx = start_x + i * (sw + gap)
    rect(s, sx, 3.05, sw, sh_box, line_color=INDIGO, lw=1.5)
    rect(s, sx, 3.05, sw, 0.42, fill=INDIGO)
    txt(s, n, sx, 3.07, sw, 0.38, size=14, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)
    txt(s, label, sx + 0.1, 3.62, sw - 0.2, 0.85, size=13, color=DARK,
        align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        ax = sx + sw + gap/2 - 0.12
        txt(s, "→", ax, 3.65, 0.3, 0.5, size=14, color=INDIGO,
            align=PP_ALIGN.CENTER)

rect(s, 0.5, 5.0, 12.3, 0.92, line_color=BORDER)
rect(s, 0.5, 5.0, 0.06, 0.92, fill=GREEN)
txt(s, "결과  ·  매수 근거와 실제 손익이 나란히 표시 → 패턴 파악 → 다음 투자 판단에 반영",
    0.72, 5.2, 11.8, 0.5, size=14, color=DARK)

txt(s, "보조 시나리오: 월말 복기(손익·승률 확인) / 대시보드에서 실시간 자산 현황 조회",
    0.5, 6.2, 12.5, 0.4, size=11, color=GRAY)


# ════════════════════════════════════════════════════════════
# 4. 기술 스택과 아키텍처
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "4.  기술 스택과 아키텍처", "주요 기술 결정 — ADR 기반")
pagenum(s, 5)

# 표 헤더
col_x = [0.5, 4.2, 9.6]
col_w = [3.6, 5.3, 3.3]
col_labels = ["영역", "선택", "근거 (ADR)"]
for lx, w, label in zip(col_x, col_w, col_labels):
    rect(s, lx, 1.15, w, 0.58, fill=INDIGO)
    txt(s, label, lx + 0.18, 1.22, w - 0.2, 0.44,
        size=13, bold=True, color=WHITE)

# 표 데이터
rows_data = [
    ("모바일 프레임워크",  "Flutter",                             "ADR-0001"),
    ("아키텍처 패턴",      "4레이어 구조\n(화면·흐름·규칙·저장)",  "ADR-0002"),
    ("백엔드 서비스",      "Firebase\n(Auth + Firestore)",        "ADR-0003"),
]
for i, (area, choice, adr) in enumerate(rows_data):
    ty = 1.73 + i * 1.52
    row_h = 1.3
    bg = LGRAY if i % 2 == 0 else WHITE

    rect(s, 0.5, ty, 3.6, row_h, fill=bg, line_color=BORDER)
    txt(s, area, 0.7, ty + 0.38, 3.2, 0.55, size=14, bold=True, color=DARK)

    rect(s, 4.2, ty, 5.3, row_h, fill=bg, line_color=BORDER)
    txt(s, choice, 4.4, ty + 0.22, 4.9, 0.9, size=14, color=DARK)

    rect(s, 9.6, ty, 3.3, row_h, fill=INDIGO, line_color=BORDER)
    txt(s, adr, 9.8, ty + 0.38, 2.9, 0.55, size=14, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)

txt(s, "* ADR 세부 내용(선택지 비교)은 발표 마지막 슬라이드에 수록",
    0.5, 6.6, 12.5, 0.45, size=11, color=GRAY)


# ════════════════════════════════════════════════════════════
# 5. 아키텍처 다이어그램
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "5.  아키텍처 다이어그램", "4레이어 구조 — 역할 분리로 유지보수 용이")
pagenum(s, 6)

layers = [
    ("Presentation",  "screens/",   "화면 렌더링 · 사용자 입력 처리",         "Flutter Widget"),
    ("Application",   "services/",  "Firebase 호출 · 상태 관리",              "FirestoreService"),
    ("Domain",        "models/",    "수익률 · 잔액 등 계산식 정의",            "Position / Account"),
    ("Data",          "Firebase",   "Firestore CRUD · Auth 인증",              "Cloud Firestore / Auth"),
]

box_h = 1.18
arrow_h = 0.18
start_y = 1.1

for i, (layer, loc, desc, tech) in enumerate(layers):
    ty = start_y + i * (box_h + arrow_h)

    rect(s, 0.5, ty, 12.3, box_h, line_color=INDIGO, lw=1.5)
    rect(s, 0.5, ty, 2.3, box_h, fill=INDIGO)
    txt(s, layer, 0.65, ty + 0.15, 2.05, 0.48, size=13, bold=True, color=WHITE)
    txt(s, loc,   0.65, ty + 0.68, 2.05, 0.38, size=10, color=RGBColor(0xC5, 0xCA, 0xE9))

    txt(s, desc, 2.95, ty + 0.2, 6.5, 0.48, size=14, bold=True, color=DARK)
    txt(s, tech, 2.95, ty + 0.7, 6.5, 0.38, size=11, color=GRAY)

    rect(s, 9.6, ty, 3.1, box_h, fill=LGRAY, line_color=BORDER)
    txt(s, tech, 9.75, ty + 0.38, 2.8, 0.42, size=11, color=GRAY)

    if i < len(layers) - 1:
        arr_x = 6.4
        txt(s, "▼", arr_x, ty + box_h - 0.02, 0.6, arrow_h + 0.05,
            size=11, color=INDIGO, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# 6. 데모
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "6.  데모", "라이브 시연 — 핵심 화면 2가지")
pagenum(s, 7)

demo_items = [
    ("시나리오 A", "계좌 관리 흐름",
     "① 계좌 등록 (주식·코인 선택, 통화 선택)  "
     "→ ② 입금 기록  "
     "→ ③ 잔액 자동 계산 확인"),
    ("시나리오 B", "포지션 기록 + 복기",
     "① 포지션 진입 기록 (종목·진입가·수량·매수 근거)  "
     "→ ② 포지션 정리 (매도가·손익·수익률 자동 계산)  "
     "→ ③ 복기 메모 조회"),
]
for i, (label, title, steps_txt) in enumerate(demo_items):
    ty = 1.2 + i * 2.5
    rect(s, 0.5, ty, 12.3, 2.1, line_color=BORDER)
    rect(s, 0.5, ty, 0.06, 2.1, fill=INDIGO)
    rect(s, 0.5, ty, 12.3, 0.52, fill=LGRAY)
    txt(s, label, 0.72, ty + 0.1, 2.0, 0.35, size=11, bold=True, color=INDIGO)
    txt(s, title, 2.85, ty + 0.1, 9.8, 0.35, size=14, bold=True, color=DARK)
    txt(s, steps_txt, 0.72, ty + 0.7, 11.9, 1.2, size=13, color=DARK)

rect(s, 0.5, 6.35, 12.3, 0.72, fill=LGRAY, line_color=BORDER)
txt(s, "데모 백업  ·  Wi-Fi 없는 환경 대비 → 로컬 데이터 캐시 + 화면 녹화 영상 준비",
    0.72, 6.53, 12.0, 0.45, size=12, color=GRAY)


# ════════════════════════════════════════════════════════════
# 7. 포트폴리오 분석 차트
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "7.  포트폴리오 분석 화면", "종목별 수익률 & 누적 실현 손익 — 앱 차트 기능")
pagenum(s, 8)

_chart_img = r"C:\project_c\portfolio-tracker\presentation\chart_screenshot.png"
if os.path.exists(_chart_img):
    _img_w = Inches(9)
    _img_x = Inches((13.33 - 9) / 2)
    s.shapes.add_picture(_chart_img, _img_x, Inches(1.3), width=_img_w)
else:
    # 스크린샷 없을 경우 matplotlib 차트 생성
    plt.rcParams['font.family'] = 'Malgun Gothic'
    plt.rcParams['axes.unicode_minus'] = False
    symbols   = ['NVDA', 'TSMC', 'AAPL', '삼성전자', 'BTC', 'ETH', 'KAKAO', 'LG엔솔']
    rates     = [35.1, 22.7, 15.2, 8.5, -12.3, 5.8, -8.9, 4.2]
    pnl_accum = [0, 1_820_000, 3_640_000, 4_560_000, 5_220_000,
                 4_580_000, 4_910_000, 4_490_000, 4_702_000]
    colors_bar = ['#1565C0' if r >= 0 else '#C62828' for r in rates]
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(13, 4.2), facecolor='#FAFAFA')
    fig.subplots_adjust(left=0.06, right=0.97, top=0.88, bottom=0.18, wspace=0.32)
    bars = ax1.bar(symbols, rates, color=colors_bar, width=0.55, edgecolor='white', linewidth=0.5)
    ax1.axhline(0, color='#9E9E9E', linewidth=0.8)
    ax1.set_title('종목별 총 수익률 (%)', fontsize=12, fontweight='bold', color='#212121', pad=8)
    ax1.set_ylabel('%', fontsize=10, color='#757575')
    ax1.tick_params(axis='x', labelsize=9, rotation=15)
    ax1.tick_params(axis='y', labelsize=9)
    ax1.set_facecolor('#FAFAFA')
    ax1.spines['top'].set_visible(False); ax1.spines['right'].set_visible(False)
    ax1.grid(axis='y', color='#E0E0E0', linewidth=0.5)
    for bar, rate in zip(bars, rates):
        ax1.text(bar.get_x() + bar.get_width()/2, bar.get_height() + (0.5 if rate >= 0 else -1.8),
                 f'{rate:+.1f}%', ha='center', va='bottom', fontsize=8, fontweight='bold',
                 color='#1565C0' if rate >= 0 else '#C62828')
    x_pts = list(range(len(pnl_accum)))
    ax2.plot(x_pts, [v/10000 for v in pnl_accum], color='#1565C0', linewidth=2.2,
             marker='o', markersize=5, markerfacecolor='#1565C0')
    ax2.fill_between(x_pts, [v/10000 for v in pnl_accum], alpha=0.12, color='#1565C0')
    ax2.set_title('누적 실현 손익 (만원)', fontsize=12, fontweight='bold', color='#212121', pad=8)
    ax2.set_ylabel('만원', fontsize=10, color='#757575')
    ax2.set_xticks(x_pts); ax2.set_xticklabels(['시작'] + symbols, fontsize=8, rotation=15)
    ax2.tick_params(axis='y', labelsize=9)
    ax2.set_facecolor('#FAFAFA')
    ax2.spines['top'].set_visible(False); ax2.spines['right'].set_visible(False)
    ax2.grid(axis='y', color='#E0E0E0', linewidth=0.5)
    final = pnl_accum[-1]/10000
    ax2.annotate(f'+{final:.0f}만원', xy=(len(pnl_accum)-1, final),
                 xytext=(-30, 12), textcoords='offset points', fontsize=10,
                 fontweight='bold', color='#1565C0',
                 arrowprops=dict(arrowstyle='->', color='#1565C0', lw=1.2))
    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='#FAFAFA')
    plt.close(fig); buf.seek(0)
    s.shapes.add_picture(buf, Inches((13.33 - 9) / 2), Inches(1.3), width=Inches(9))

rect(s, 0.4, 6.55, 12.5, 0.62, fill=LGRAY, line_color=BORDER)
txt(s, "기간 필터(1주·1달·3달·6달·1년·전체) 적용 시 해당 기간 정리 포지션만 반영  "
       "·  분할 매도 시 동일 종목 합산 표시",
    0.6, 6.68, 12.1, 0.4, size=11, color=GRAY)


# ════════════════════════════════════════════════════════════
# 8. 개발 환경 · 빌드 & 배포
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "8.  개발 환경 · 빌드 & 배포", "설치부터 배포까지 — GitHub 가이드 참고")
pagenum(s, 9)

env_items = [
    ("개발 환경 구성",
     "Flutter SDK 3.x  ·  Dart SDK 3.x  ·  VS Code + Flutter/Dart 확장\n"
     "Firebase CLI 15.x  ·  Android Studio (에뮬레이터·SDK)  ·  Git\n"
     "설치 가이드: README.md → docs/setup.md 순서대로 진행"),
    ("빌드",
     "flutter pub get  →  flutter run  (개발·디버그 빌드)\n"
     "flutter build apk --release  →  APK 생성  ·  build/app/outputs/apk/release/"),
    ("배포 & 설치",
     "Firebase Auth + Firestore: 이미 연동 완료 (firebase_options.dart)\n"
     "Android APK 직접 설치: USB 케이블 + adb install app-release.apk\n"
     "Firebase Hosting (web): flutter build web → firebase deploy"),
]
for i, (title, body) in enumerate(env_items):
    ty = 1.15 + i * 1.82
    rect(s, 0.5, ty, 12.3, 1.55, line_color=BORDER)
    rect(s, 0.5, ty, 0.06, 1.55, fill=INDIGO)
    rect(s, 0.5, ty, 12.3, 0.48, fill=LGRAY)
    txt(s, title, 0.72, ty + 0.1, 12.0, 0.3, size=13, bold=True, color=INDIGO)
    txt(s, body,  0.72, ty + 0.56, 11.8, 0.88, size=11.5, color=DARK)


# ════════════════════════════════════════════════════════════
# 10. 코드 품질 · 테스트 · 성능 최적화
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "9.  코드 품질 · 테스트 · 성능 최적화")
pagenum(s, 10)

quality_items = [
    ("코드 품질 관리",
     "flutter_lints 적용 — 린트 경고 0건 유지  ·  4레이어 아키텍처로 역할 분리\n"
     "FirestoreService 단일 책임 원칙 준수  ·  모델·서비스·화면 계층 엄격히 분리"),
    ("단위 테스트 · 통합 테스트",
     "docs/testing.md 기반 기능별 수동 체크리스트 작성 및 확인\n"
     "인증·계좌·포지션·차트·대시보드 5개 영역 E2E 시나리오 직접 수행\n"
     "Firebase Auth + Firestore 실제 연동 통합 테스트 완료"),
    ("성능 최적화",
     "StreamBuilder — Firestore 실시간 구독으로 불필요한 poll 제거\n"
     "계좌 삭제 시 Batch 처리로 하위 트랜잭션 원자적 삭제 (cascade)\n"
     "FutureBuilder 캐싱으로 계좌 목록 잔액 계산 중복 호출 최소화"),
]
for i, (title, body) in enumerate(quality_items):
    ty = 1.15 + i * 1.82
    rect(s, 0.5, ty, 12.3, 1.55, line_color=BORDER)
    rect(s, 0.5, ty, 0.06, 1.55, fill=AMBER)
    rect(s, 0.5, ty, 12.3, 0.48, fill=LGRAY)
    txt(s, title, 0.72, ty + 0.1, 12.0, 0.3, size=13, bold=True, color=INDIGO)
    txt(s, body,  0.72, ty + 0.56, 11.8, 0.88, size=11.5, color=DARK)


# ════════════════════════════════════════════════════════════
# 11. 어려운 점 / 도움 요청
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "10.  어려운 점 / 도움 요청")
pagenum(s, 11)

issues = [
    ("iOS 빌드 불가",
     "Mac 환경 없음 → iOS 시뮬레이터·실기기 테스트 불가\nAndroid APK만 제출 예정"),
    ("Firebase Security Rules",
     "Firestore 보안 규칙 검증 어려움 — 다른 사용자가 타 계정 데이터에\n"
     "접근하지 못하도록 규칙을 작성했으나 침투 테스트 환경 없음"),
    ("앱 배포 경험 없음",
     "Google Play Console 등록 절차 및 APK 서명 과정 첫 시도\n"
     "배포 오류 발생 시 대응 방법 조언 필요"),
]
for i, (title, body) in enumerate(issues):
    ty = 1.15 + i * 1.82
    rect(s, 0.5, ty, 12.3, 1.55, line_color=BORDER)
    rect(s, 0.5, ty, 0.06, 1.55, fill=AMBER)
    txt(s, title, 0.72, ty + 0.15, 12.0, 0.45, size=14, bold=True, color=DARK)
    txt(s, body,  0.72, ty + 0.62, 11.8, 0.82, size=13, color=DARK)


# ════════════════════════════════════════════════════════════
# 9. 개발 방식 — AI 워크플로우 (가산점 A+1 · B+2 · C+1)
# ════════════════════════════════════════════════════════════
# 9. 개발 방식 — AI 워크플로우 (주성현 1/2)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "11.  개발 방식 — AI 워크플로우", "주성현  (1/2)")
pagenum(s, 12)

def ai_item(s, ty, bh, badge_color, badge_txt, title, body):
    rect(s, 0.5, ty, 12.3, bh, line_color=BORDER)
    rect(s, 0.5, ty, 12.3, 0.44, fill=LGRAY)
    rect(s, 0.5, ty, 1.6, 0.44, fill=badge_color)
    txt(s, badge_txt, 0.5, ty, 1.6, 0.44,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, 2.25, ty + 0.07, 10.4, 0.30, size=12, bold=True, color=INDIGO)
    txt(s, body,  0.72, ty + 0.52, 11.8, bh - 0.58, size=10.5, color=DARK)

for i, (badge_color, badge_txt, title, body) in enumerate([
    (AMBER,  "AI 워크플로우\n활용  +1",
     "feature-debugger 에이전트 직접 제작 + /dg 실사용",
     ".claude/agents/feature-debugger.md  직접 제작\n"
     "앱 기능 6개 영역(인증 · 계좌 · 포지션 · 차트 · 대시보드 · Firebase)을 /dg 명령어 하나로 반복 점검 수행"),
    (INDIGO, "본인 기법\n구성  +2",
     "기법 1  —  단일 md 부트스트랩  (AUTHORING.Jooseonghyeon.v0.1.0.md)",
     "Rules · Skills · Commands 전체를 파일 하나로 자동 세팅\n"
     "새 프로젝트에 이 파일 하나만 주면 AI 운영 체계 즉시 구성"),
    (INDIGO, "본인 기법\n구성  +2",
     "기법 3  —  단계별 개발 계획 · 진행",
     "개발 전 단계별 계획 수립 → AI 구현 → 직접 점검 · 이해 → 다음 단계 진행\n"
     "AI 의존 개발에서 코드 이해도를 유지하기 위한 개인 워크플로우"),
]):
    ai_item(s, 1.0 + i * 2.0, 1.85, badge_color, badge_txt, title, body)


# ════════════════════════════════════════════════════════════
# 10. 개발 방식 — AI 워크플로우 (주성현 2/2)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "8.  개발 방식 — AI 워크플로우", "주성현  (2/2)")
pagenum(s, 13)

for i, (badge_color, badge_txt, title, body) in enumerate([
    (GREEN,  "LLM Wiki\n암묵지  +1",
     "LLM Wiki  —  암묵지 운영  (notes/llm-wiki.md · 11개 항목)",
     "개발 중 배운 AI 사용 인사이트 직접 기록\n"
     "Flutter 실행 환경 / 에이전트 호출법 / ADR 개념 / python-pptx 단위 변환 등"),
    (RGBColor(0x6A, 0x1B, 0x9A), "AI 리포트\n발표  +2",
     "AI Agent 리포트  —  Claude Code 활용 사례 발표",
     "Claude Code(Anthropic, 2024~)로 서브에이전트 제작 · 커스텀 명령어 운영 · 단계별 개발 수행\n"
     "최근 6개월 이내 지속 업데이트 도구 · 출처: claude.ai/code  (공식 문서)"),
]):
    ai_item(s, 0.9 + i * 2.95, 2.7, badge_color, badge_txt, title, body)


# ════════════════════════════════════════════════════════════
# 11. 개발 방식 — AI 워크플로우 (정경우 1/2)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "9.  개발 방식 — AI 워크플로우", "정경우  (1/2)")
pagenum(s, 14)

for i, (badge_color, badge_txt, title, body) in enumerate([
    (AMBER,  "AI 워크플로우\n활용  +1",
     "pre-presentation-checker 에이전트 직접 제작 + /pc 실사용",
     ".claude/agents/pre-presentation-checker.md  직접 제작\n"
     "발표 전 제출 서류 · 가산점 · 앱 상태를 평가 루브릭 기준으로 자동 점검 → 미완료 항목 우선 보고"),
    (INDIGO, "본인 기법\n구성  +2",
     "기법 1  —  단일 md 부트스트랩  (AUTHORING.jeongkyeongwoo.v0.1.0.md)",
     "Rules · Skills · Commands 전체를 파일 하나로 자동 세팅\n"
     "새 프로젝트에 이 파일 하나만 주면 AI 운영 체계 즉시 구성"),
    (INDIGO, "본인 기법\n구성  +2",
     "기법 3  —  단계별 보고 방식",
     "각 구현 단계 완료 후 AI에게 결과 보고받기 → 코드 이해도 유지\n"
     "\"방금 만든 기능을 3줄로 설명해줘\" → 내용 이해 후 다음 단계 진행  ·  Q&A 즉답 가능"),
]):
    ai_item(s, 1.0 + i * 2.0, 1.85, badge_color, badge_txt, title, body)


# ════════════════════════════════════════════════════════════
# 12. 개발 방식 — AI 워크플로우 (정경우 2/2)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "9.  개발 방식 — AI 워크플로우", "정경우  (2/2)")
pagenum(s, 15)

for i, (badge_color, badge_txt, title, body) in enumerate([
    (GREEN,  "LLM Wiki\n암묵지  +1",
     "LLM Wiki  —  암묵지 운영  (notes/llm-wiki.md · 11개 항목)",
     "개발 중 배운 AI 사용 인사이트 직접 기록\n"
     "Flutter 실행 환경 / 에이전트 호출법 / ADR 개념 / python-pptx 단위 변환 등"),
    (RGBColor(0x6A, 0x1B, 0x9A), "AI 리포트\n발표  +2",
     "AI Agent 리포트  —  Claude Code 활용 사례 발표",
     "Claude Code(Anthropic, 2024~)로 서브에이전트 제작 · 커스텀 명령어 운영 · 단계별 보고 방식 수행\n"
     "최근 6개월 이내 지속 업데이트 도구 · 출처: claude.ai/code  (공식 문서)"),
]):
    ai_item(s, 0.9 + i * 2.95, 2.7, badge_color, badge_txt, title, body)


# ════════════════════════════════════════════════════════════
# 11. ADR-0001  모바일 프레임워크
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "ADR-0001  —  모바일 프레임워크 선택")
pagenum(s, 16)

txt(s, "선택지 비교", 0.5, 1.08, 12.0, 0.38, size=12, color=GRAY)
for i, (opt, sel) in enumerate([
    ("React Native", False),
    ("안드로이드 전용 개발", False),
    ("아이폰 전용 개발", False),
    ("Flutter", True),
]):
    chip(s, opt, 0.5, 1.52 + i * 1.38, 12.3, 1.12, selected=sel)


# ════════════════════════════════════════════════════════════
# 12. ADR-0002  아키텍처 패턴
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "ADR-0002  —  아키텍처 패턴 선택")
pagenum(s, 17)

txt(s, "선택지 비교", 0.5, 1.08, 6.0, 0.38, size=12, color=GRAY)
for i, (opt, sel) in enumerate([
    ("역할 구분 없이 한 파일", False),
    ("MVC 방식", False),
    ("6 ~ 8개 레이어", False),
    ("4레이어 구조", True),
]):
    chip(s, opt, 0.5, 1.52 + i * 1.38, 6.0, 1.12, selected=sel)

# 4레이어 표 (우측)
txt(s, "선택된 4레이어 구조", 7.0, 1.08, 5.9, 0.38, size=12, color=GRAY)
table_rows = [
    ("레이어",  "역할",                   "위치",        True),
    ("화면",    "버튼·목록·입력창",        "screens/",    False),
    ("흐름",    "화면-데이터 연결 처리",   "",            False),
    ("규칙",    "잔액·수익률 계산",        "models/",     False),
    ("저장",    "Firestore 입출력",        "services/",   False),
]
for i, (a, b, c, hdr) in enumerate(table_rows):
    ty = 1.52 + i * 1.1
    bg = LGRAY if hdr else WHITE
    for val, lx, w in [(a, 7.0, 1.5), (b, 8.55, 2.8), (c, 11.4, 1.9)]:
        rect(s, lx, ty, w, 0.9, fill=bg, line_color=BORDER)
        txt(s, val, lx + 0.12, ty + 0.22, w - 0.2, 0.48,
            size=12, bold=hdr, color=INDIGO if hdr else DARK)


# ════════════════════════════════════════════════════════════
# 13. ADR-0003  백엔드 서비스
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "ADR-0003  —  백엔드 서비스 선택")
pagenum(s, 18)

txt(s, "선택지 비교", 0.5, 1.08, 12.0, 0.38, size=12, color=GRAY)
for i, (opt, sel) in enumerate([
    ("Supabase", False),
    ("폰 내부 저장 (SQLite 등)", False),
    ("직접 서버 구축", False),
    ("Firebase (Auth + Firestore)", True),
]):
    chip(s, opt, 0.5, 1.52 + i * 1.38, 12.3, 1.12, selected=sel)


# ════════════════════════════════════════════════════════════
# 감사합니다
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=WHITE)
rect(s, 0, 5.1, 13.33, 0.05, fill=INDIGO)
rect(s, 0, 5.15, 13.33, 2.35, fill=LGRAY)

txt(s, "감사합니다", 1, 1.6, 11, 1.5,
    size=52, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
txt(s, "Portfolio Tracker  ·  2026.06",
    1, 5.55, 11, 0.5, size=13, color=GRAY, align=PP_ALIGN.CENTER)
pagenum(s, 19)


# ── 저장 ─────────────────────────────────────────────────
out = r"C:\project_c\portfolio-tracker\presentation\portfolio-tracker-final.pptx"
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print(f"저장 완료: {out}")
