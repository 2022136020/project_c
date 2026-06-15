from pptx import Presentation
from pptx.util import Inches
import re

PPT_PATH   = r"C:\project_c\portfolio-tracker\presentation\portfolio-tracker-final.pptx"
TOTAL      = 23   # 새 총 슬라이드 수

# 번호 없는 슬라이드 인덱스(0-based): 표지(0), 시연영상(21), 감사합니다(22)
NO_NUM = {0, 21, 22}

prs = Presentation(PPT_PATH)

THRESHOLD = Inches(0.4)

def is_pagenum_box(shape):
    if not shape.has_text_frame:
        return False
    text = shape.text_frame.text.strip()
    # "N / M" 형식 또는 숫자만
    if not (re.match(r'^\d+\s*/\s*\d+$', text) or text.isdigit()):
        return False
    left_ok = abs(shape.left - Inches(12.3)) < THRESHOLD
    top_ok  = abs(shape.top  - Inches(7.1))  < THRESHOLD
    return left_ok and top_ok

for slide_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if not is_pagenum_box(shape):
            continue

        if slide_idx in NO_NUM:
            sp = shape._element
            sp.getparent().remove(sp)
            print(f"slide {slide_idx+1:2d}: 번호 제거")
        else:
            new_num  = slide_idx + 1
            new_text = f"{new_num} / {TOTAL}"
            tf = shape.text_frame
            # 기존 run 텍스트 교체
            old = tf.paragraphs[0].runs[0].text if tf.paragraphs[0].runs else "?"
            tf.paragraphs[0].runs[0].text = new_text
            print(f"slide {slide_idx+1:2d}: '{old}' -> '{new_text}'")

prs.save(PPT_PATH)
print("\n저장 완료")
