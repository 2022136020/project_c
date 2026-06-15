from pptx import Presentation
from pptx.util import Inches
import re

PPT_PATH = r"C:\project_c\portfolio-tracker\presentation\portfolio-tracker-final.pptx"

prs = Presentation(PPT_PATH)

# 슬라이드 인덱스(0-base) → 새 섹션 번호
# (번호 없는 슬라이드: 표지0, ADR 10~12, GitHub가이드16, 시연영상21, 감사합니다22 제외)
NEW_NUM = {
    1:  1,   # 문제 정의      (기존 1 → 그대로)
    2:  2,   # WBS           (기존 3 → 2)
    3:  3,   # 기술 스택      (기존 5 → 3)
    4:  4,   # AI 주성현 1/2  (기존 13 → 4)
    5:  5,   # AI 주성현 2/2  (기존 14 → 5)
    6:  6,   # AI 정경우 1/2  (기존 15 → 6)
    7:  7,   # AI 정경우 2/2  (기존 16 → 7)
    8:  8,   # 아키텍처        (기존 6 → 8)
    9:  9,   # 활용 방안       (기존 2 → 9)
    13: 10,  # 개발 환경       (기존 9 → 10)
    14: 11,  # 어려운 점       (기존 11 → 그대로)
    15: 12,  # 코드 품질       (기존 10 → 12)
    17: 13,  # 데모            (기존 7 → 13)
    18: 14,  # 사용자 시나리오  (기존 4 → 14)
    19: 15,  # 포트폴리오 분석  (기존 8 → 15)
    20: 16,  # 향후 계획        (기존 12 → 16)
}

TITLE_LEFT_IN = 0.5
TITLE_TOP_IN  = 0.08
THRESHOLD = Inches(0.3)

def is_title_box(shape):
    left_ok = abs(shape.left - Inches(TITLE_LEFT_IN)) < THRESHOLD
    top_ok  = abs(shape.top  - Inches(TITLE_TOP_IN))  < THRESHOLD
    return shape.has_text_frame and left_ok and top_ok

for slide_idx, slide in enumerate(prs.slides):
    new_n = NEW_NUM.get(slide_idx)
    if new_n is None:
        continue

    for shape in slide.shapes:
        if not is_title_box(shape):
            continue
        tf = shape.text_frame
        if not tf.paragraphs or not tf.paragraphs[0].runs:
            continue
        run = tf.paragraphs[0].runs[0]
        old_text = run.text

        # "13.  개발 방식 — AI 워크플로우" → "4.  개발 방식 — AI 워크플로우"
        new_text = re.sub(r'^\d+\.\s+', f"{new_n}.  ", old_text)
        if new_text != old_text:
            run.text = new_text
            print(f"slide {slide_idx+1:2d}: updated -> {new_n}.")
        else:
            print(f"slide {slide_idx+1:2d}: no change ({new_n}.)")

prs.save(PPT_PATH)
print("\n저장 완료")
