from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import qrcode, io

PPT_PATH   = r"C:\project_c\portfolio-tracker\presentation\portfolio-tracker-final.pptx"
GITHUB_URL = "https://github.com/2022136020/project_c/tree/main/portfolio-tracker"

INDIGO = RGBColor(0x3F, 0x51, 0xB5)
GRAY   = RGBColor(0x9E, 0x9E, 0x9E)
DARK   = RGBColor(0x21, 0x21, 0x21)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation(PPT_PATH)
W = prs.slide_width
H = prs.slide_height

# ─────────────────────────────────────────────────────────────
# 1. 슬라이드 재순서
#
# 현재 22슬라이드 (0-index):
#  0 표지               8 포트폴리오 분석    16 AI 워크플로우 정경우 1/2
#  1 문제 정의           9 개발환경·빌드배포   17 ADR-0001
#  2 활용방안&기대효과   10 코드품질·테스트    18 ADR-0002
#  3 WBS               11 어려운 점         19 ADR-0003
#  4 사용자 시나리오     12 향후 계획         20 시연 영상
#  5 기술 스택           13 AI 워크플로우 주성현 1/2  21 감사합니다
#  6 아키텍처 다이어그램  14 AI 워크플로우 주성현 2/2
#  7 데모               15 AI 워크플로우 정경우 1/2
#
# 교수님 순서 기준 새 순서:
#  비전제시 → 문제정의 → WBS → 기술스택 → AI워크플로우(x4)
#  → 아키텍처 → 활용방안 → ADR(x3)
#  → 개발환경/빌드배포 → 어려운점 → 코드품질 → [깃허브가이드:NEW]
#  → 데모 → 사용자시나리오 → 포트폴리오분석 → 향후계획
#  → 시연영상 → 감사합니다
# ─────────────────────────────────────────────────────────────

new_order = [0, 1, 3, 5, 13, 14, 15, 16, 6, 2, 17, 18, 19, 9, 11, 10, 7, 4, 8, 12, 20, 21]
# 깃허브 가이드는 인덱스 16 위치(코드품질 다음, 데모 앞)에 삽입

sldIdLst = prs.slides._sldIdLst
original_ids = list(sldIdLst)

for el in original_ids:
    sldIdLst.remove(el)
for idx in new_order:
    sldIdLst.append(original_ids[idx])

print(f"슬라이드 재순서 완료: {len(prs.slides)}개")

# ─────────────────────────────────────────────────────────────
# 2. GitHub 설치 가이드 슬라이드 추가 (현재 마지막에 붙은 뒤 위치 이동)
# ─────────────────────────────────────────────────────────────
BLANK = prs.slide_layouts[6]
slide = prs.slides.add_slide(BLANK)

# 헤더 (기존 PPT 스타일 동일)
def add_rect(s, l, t, w, h, fill=None):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    sh.line.fill.background()
    return sh

def add_txt(s, text, l, t, w, h, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    return tb

# 헤더 바
add_rect(slide, 0, 0, 13.33, 1.0)
add_txt(slide, "GitHub 설치 가이드", 0.5, 0.08, 12, 0.62,
        size=24, bold=True, color=INDIGO)
add_txt(slide, "포트폴리오 트래커 앱 — 클론 · 설치 · 실행 방법",
        0.5, 0.68, 12, 0.3, size=11, color=GRAY)
add_rect(slide, 0, 0.95, 13.33, 0.05, fill=INDIGO)

# 페이지 번호 (17번)
pn = slide.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.3))
r = pn.text_frame.paragraphs[0].add_run()
r.text = "17"
r.font.size = Pt(10); r.font.color.rgb = GRAY

# ── QR 코드 생성 ──
qr = qrcode.QRCode(box_size=10, border=2,
                   error_correction=qrcode.constants.ERROR_CORRECT_M)
qr.add_data(GITHUB_URL)
qr.make(fit=True)
qr_img = qr.make_image(fill_color="black", back_color="white")
buf = io.BytesIO()
qr_img.save(buf, format="PNG")
buf.seek(0)

# QR 코드 배치 (왼쪽)
slide.shapes.add_picture(buf, Inches(0.8), Inches(1.3), Inches(3.5), Inches(3.5))

# QR 아래 URL 텍스트
add_txt(slide, GITHUB_URL, 0.4, 4.9, 4.5, 0.4, size=9, color=GRAY, align=PP_ALIGN.CENTER)

# ── 설치 단계 (오른쪽) ──
steps = [
    ("1.  저장소 클론",
     "git clone https://github.com/2022136020/project_c.git\ncd project_c/portfolio-tracker"),
    ("2.  패키지 설치",
     "flutter pub get"),
    ("3.  앱 실행",
     "flutter run"),
    ("4.  APK 빌드 (Android)",
     "flutter build apk --release"),
]

y = 1.25
for title, cmd in steps:
    add_txt(slide, title, 5.0, y, 7.8, 0.35, size=13, bold=True, color=INDIGO)
    y += 0.35
    # 코드 배경 박스
    box_h = 0.28 * (cmd.count('\n') + 1) + 0.1
    add_rect(slide, 5.0, y, 7.8, box_h, fill=RGBColor(0xF5, 0xF5, 0xF5))
    add_txt(slide, cmd, 5.15, y + 0.02, 7.5, box_h,
            size=10, color=DARK)
    y += box_h + 0.2

# 하단 안내
add_txt(slide,
        "상세 가이드: docs/setup.md  ·  docs/deploy.md  ·  portfolio-tracker/README.md",
        0.5, 6.85, 12.3, 0.4, size=10, color=GRAY, align=PP_ALIGN.CENTER)

# ── 방금 추가한 슬라이드를 인덱스 16 위치로 이동 ──
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
last = ids[-1]
sldIdLst.remove(last)
# lxml insert
sldIdLst.insert(16, last)

print(f"GitHub 가이드 슬라이드 추가 완료 (위치: 17번)")
print(f"최종 슬라이드 수: {len(prs.slides)}")

prs.save(PPT_PATH)
print("저장 완료")
