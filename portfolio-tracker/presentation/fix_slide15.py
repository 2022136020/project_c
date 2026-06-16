# -*- coding: utf-8 -*-
from pptx import Presentation
from lxml import etree

PPT_PATH = r"C:\project_c\portfolio-tracker\presentation\portfolio-tracker-final.pptx"
prs = Presentation(PPT_PATH)

slide = prs.slides[14]  # 슬라이드 15 (0-indexed)

ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
changed = False

for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    for t_el in shape._element.iter(f"{{{ns}}}t"):
        if "배포 오류 발생 시" in t_el.text:
            print(f"Before: {repr(t_el.text)}")
            t_el.text = t_el.text.replace("배포 오류 발생 시", "배포 및 오류 발생 시")
            print(f"After:  {repr(t_el.text)}")
            changed = True

if changed:
    prs.save(PPT_PATH)
    print("저장 완료")
else:
    print("변경 대상 없음")
